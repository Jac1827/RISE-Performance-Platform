import csv
import subprocess
import sys
import zipfile
from collections import defaultdict
from datetime import datetime
from pathlib import Path
from xml.etree import ElementTree as ET

sys.path.insert(0, "/tmp/pptxlib")

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path("/Users/jacheflin/Desktop/55+ 90 DAY PLAN SD")
OUTPUT = Path("/Users/jacheflin/Documents/Playground/RISE_GKP_Viera_Q1_2026_Review.pptx")
LOGO_SOURCE = Path("/Users/jacheflin/Desktop/LOGO.exr")
LOGO_PNG = Path("/tmp/rise_logo_brand.png")
LOGO_SAFE = Path("/tmp/rise_logo_brand_safe.png")

BOX_2026_AVAIL = BASE / "02 - RISE - Box Score - Availability Q12026).csv"
BOX_2026_LEADS = BASE / "02 - RISE - Box Score - Lead Activity (01-01-2026 - 03-31-2026).csv"
BOX_2026_CONV = BASE / "02 - RISE - Box Score - Lead Conversions (01-01-2026 - 03-31-2026).csv"
BOX_2026_PULSE = BASE / "02 - RISE - Box Score - Property Pulse (01-01-2026 - 03-31-2026).csv"

BOX_2025_DIR = BASE / "02 - RISE - Box Score - 03-24-2026"
BOX_2025_AVAIL = BOX_2025_DIR / "02 - RISE - Box Score - Availability (As of 01-31-2025).csv"
BOX_2025_LEADS = BOX_2025_DIR / "02 - RISE - Box Score - Lead Activity (01-01-2025 - 01-31-2025).csv"
BOX_2025_CONV = BOX_2025_DIR / "02 - RISE - Box Score - Lead Conversions (01-01-2025 - 01-31-2025).csv"
BOX_2025_PULSE = BOX_2025_DIR / "02 - RISE - Box Score - Property Pulse (01-01-2025 - 01-31-2025).csv"

INCOME_2026 = BASE / "RISE - Budget Comparison - Income Statement.xlsx"
INCOME_2025 = BASE / "RISE - Budget Comparison - Income Statement-2.xlsx"

GKP_TYPES = {"A1C", "A2E", "A2G", "A2H", "A3F", "A4C", "B1A", "B1K", "B3A", "H2D", "H2F"}
VIERA_TYPES = {"A1", "A2", "A3", "A4-A", "A4-B", "B1", "B3", "H1C", "H2A", "H3B"}

AVAIL_COLS = {
    "budget_rent": 2,
    "scheduled_rent": 3,
    "net_eff": 4,
    "units": 5,
    "occupied_units": 8,
    "available": 10,
    "notice_rented": 12,
    "vacant_rented": 14,
    "physical_occ_rate": 16,
    "leased_rate": 17,
}

NAVY = RGBColor(2, 48, 70)
BLUE = RGBColor(84, 157, 184)
PAPER = RGBColor(248, 252, 253)
CARD = RGBColor(237, 246, 249)
INK = RGBColor(2, 48, 70)
MUTED = RGBColor(79, 118, 134)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(24, 156, 83)
AMBER = RGBColor(230, 146, 0)
RED = RGBColor(198, 82, 67)
BORDER = RGBColor(210, 226, 233)

TRAINING_SUMMARY = (
    "Priority 1 training focus: fix lead intake and first-response reliability, improve discovery and "
    "follow-up discipline, and protect closing ratios."
)
VIERA_TRAINING_DATE = "Wednesday 04/15/26"
GKP_TRAINING_DATE = "Thursday 04/23/26"
VIERA_FOLLOWUP_DATE = "Wednesday 05/13/26"
GKP_FOLLOWUP_DATE = "Thursday 05/21/26"

FILTERED_ABSORPTION = {
    "Viera": {
        "market": "Viera active-adult benchmark",
        "subject_name": "RISE Viera",
        "subject_units": 166,
        "subject_leases30": 7,
        "subject_abs": 0.0422,
        "subject_leased": 0.72,
        "submarket_abs": 0.0189,
        "submarket_leased": 0.85,
        "delta_abs": 0.0233,
        "delta_leased": -0.13,
        "comps": ["Parasol Melbourne"],
        "status": "Ahead on pace",
        "status_color": GREEN,
        "growth_note": "Filtered pace is modest and benchmark depth is thin: the shared report includes one clean active-adult comp.",
    },
    "GKP": {
        "market": "Jacksonville active-adult benchmark",
        "subject_name": "RISE at Glen Kernan Park",
        "subject_units": 308,
        "subject_leases30": 9,
        "subject_abs": 0.0292,
        "subject_leased": 0.13,
        "submarket_abs": 0.0367,
        "submarket_leased": 0.8393,
        "delta_abs": -0.0075,
        "delta_leased": -0.7093,
        "comps": ["Olea at Beach Haven", "Olea eTown", "RISE at Nocatee"],
        "status": "Slightly slower",
        "status_color": AMBER,
        "growth_note": "Filtered Jacksonville 55+ / active-adult peers are leasing at a healthier clip, led by Olea eTown.",
    },
}


def num(value):
    if value is None or value == "":
        return 0.0
    return float(str(value).replace(",", ""))


def num_i(value):
    return int(round(num(value)))


def pct(value, digits=1):
    return f"{value * 100:.{digits}f}%"


def money(value):
    return "${:,.0f}".format(value)


def delta_pct(actual, baseline, digits=1):
    if baseline == 0:
        return None
    return f"{(actual / baseline - 1) * 100:+.{digits}f}%"


def delta_pts(actual, baseline, digits=1):
    return f"{(actual - baseline) * 100:+.{digits}f} pts"


def read_rows(path):
    with path.open(newline="", encoding="utf-8-sig") as fh:
        return list(csv.reader(fh))


def ensure_logo_png():
    if not LOGO_SOURCE.exists():
        return None
    if LOGO_SAFE.exists() and LOGO_SAFE.stat().st_mtime >= LOGO_SOURCE.stat().st_mtime:
        return LOGO_SAFE
    subprocess.run(
        ["sips", "-s", "format", "png", str(LOGO_SOURCE), "--out", str(LOGO_PNG)],
        check=True,
        capture_output=True,
        text=True,
    )
    with Image.open(LOGO_PNG) as img:
        # PowerPoint is more reliable with a standard 8-bit RGB asset than a 16-bit RGBA file from EXR conversion.
        base = Image.new("RGB", img.size, (255, 255, 255))
        if img.mode in {"RGBA", "LA"}:
            base.paste(img.convert("RGBA"), mask=img.convert("RGBA").split()[-1])
        else:
            base.paste(img.convert("RGB"))
        base.save(LOGO_SAFE, format="PNG", optimize=True)
    return LOGO_SAFE


def parse_availability(path):
    rows = read_rows(path)[1:]
    output = {}
    for name, unit_types in [("GKP", GKP_TYPES), ("Viera", VIERA_TYPES)]:
        subset = [row for row in rows if row[0] in unit_types]
        units = sum(num_i(row[AVAIL_COLS["units"]]) for row in subset)
        occupied_units = sum(num_i(row[AVAIL_COLS["occupied_units"]]) for row in subset)
        physical_occ = sum(
            num(row[AVAIL_COLS["physical_occ_rate"]]) * num_i(row[AVAIL_COLS["units"]])
            for row in subset
        ) / units
        leased = sum(
            num(row[AVAIL_COLS["leased_rate"]]) * num_i(row[AVAIL_COLS["units"]]) for row in subset
        ) / units
        avg_net_eff = sum(
            num(row[AVAIL_COLS["net_eff"]]) * num_i(row[AVAIL_COLS["occupied_units"]])
            for row in subset
            if num_i(row[AVAIL_COLS["occupied_units"]]) > 0
        ) / max(occupied_units, 1)
        avg_sched = sum(
            num(row[AVAIL_COLS["scheduled_rent"]]) * num_i(row[AVAIL_COLS["occupied_units"]])
            for row in subset
            if num_i(row[AVAIL_COLS["occupied_units"]]) > 0
        ) / max(occupied_units, 1)
        output[name] = {
            "units": units,
            "occupied_units": occupied_units,
            "physical_occ": physical_occ,
            "leased": leased,
            "avg_net_eff": avg_net_eff,
            "avg_sched": avg_sched,
            "available": sum(num_i(row[AVAIL_COLS["available"]]) for row in subset),
            "vacant_rented": sum(num_i(row[AVAIL_COLS["vacant_rented"]]) for row in subset),
            "notice_rented": sum(num_i(row[AVAIL_COLS["notice_rented"]]) for row in subset),
        }
    return output


def parse_lead_activity(path):
    rows = read_rows(path)[1:]
    totals = [row for row in rows if row[0] == "Unknown"]
    return {
        "GKP": {"leads": num_i(totals[0][1]), "tours": num_i(totals[0][14])},
        "Viera": {"leads": num_i(totals[1][1]), "tours": num_i(totals[1][14])},
    }


def parse_lead_conversions(path):
    rows = read_rows(path)[1:]
    output = {"GKP": defaultdict(int), "Viera": defaultdict(int)}
    for row in rows:
        name = "GKP" if row[0] in GKP_TYPES else "Viera"
        output[name]["apps_completed"] += num_i(row[1])
        output[name]["apps_approved"] += num_i(row[5])
        output[name]["leases_completed"] += num_i(row[7])
        output[name]["leases_approved"] += num_i(row[9])
    return output


def parse_property_pulse(path):
    rows = read_rows(path)[1:]
    output = {"GKP": defaultdict(int), "Viera": defaultdict(int)}
    for row in rows:
        name = "GKP" if row[0] in GKP_TYPES else "Viera"
        output[name]["move_ins"] += num_i(row[2])
        output[name]["notices"] += num_i(row[7])
        output[name]["move_outs"] += num_i(row[8])
        output[name]["renewal_approvals"] += num_i(row[5])
    return output


def read_xlsx_sheet(path, sheet_name):
    ns = {"a": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
    with zipfile.ZipFile(path) as archive:
        shared_strings = []
        if "xl/sharedStrings.xml" in archive.namelist():
            root = ET.fromstring(archive.read("xl/sharedStrings.xml"))
            for item in root.findall("a:si", ns):
                shared_strings.append("".join(node.text or "" for node in item.iterfind(".//a:t", ns)))

        workbook = ET.fromstring(archive.read("xl/workbook.xml"))
        rels = ET.fromstring(archive.read("xl/_rels/workbook.xml.rels"))
        rel_map = {rel.attrib["Id"]: rel.attrib["Target"] for rel in rels}

        target = None
        for sheet in workbook.find("a:sheets", ns):
            if sheet.attrib["name"] == sheet_name:
                rid = sheet.attrib["{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"]
                target = "xl/" + rel_map[rid]
                break

        root = ET.fromstring(archive.read(target))
        rows = []
        for row in root.find("a:sheetData", ns):
            values = {}
            for cell in row:
                ref = cell.attrib["r"]
                col = "".join(ch for ch in ref if ch.isalpha())
                value_node = cell.find("a:v", ns)
                if value_node is None:
                    value = ""
                else:
                    value = value_node.text or ""
                    if cell.attrib.get("t") == "s":
                        value = shared_strings[int(value)]
                values[col] = value
            rows.append(values)
        return rows


def parse_income_statement(path):
    output = {}
    for name, sheet_name in [("GKP", "Rise At Glen Kernan Park"), ("Viera", "Rise Viera")]:
        rows = read_xlsx_sheet(path, sheet_name)
        metrics = {}
        for row in rows:
            label = row.get("B", "")
            if label in {"Gross Potential Rent (GPR)", "Net Rental Income", "Total Income", "Total Cost Of Leasing"}:
                metrics[label] = {
                    "actual": num(row.get("C", "")),
                    "budget": num(row.get("D", "")),
                    "variance": num(row.get("E", "")),
                    "variance_pct": num(row.get("F", "")),
                }
        output[name] = metrics
    return output


def assemble_data():
    avail_2026 = parse_availability(BOX_2026_AVAIL)
    leads_2026 = parse_lead_activity(BOX_2026_LEADS)
    conv_2026 = parse_lead_conversions(BOX_2026_CONV)
    pulse_2026 = parse_property_pulse(BOX_2026_PULSE)

    avail_2025 = parse_availability(BOX_2025_AVAIL)
    leads_2025 = parse_lead_activity(BOX_2025_LEADS)
    conv_2025 = parse_lead_conversions(BOX_2025_CONV)
    pulse_2025 = parse_property_pulse(BOX_2025_PULSE)

    income_2026 = parse_income_statement(INCOME_2026)
    income_2025 = parse_income_statement(INCOME_2025)

    submarket_avg_rent = {"GKP": 2231, "Viera": 2111}
    reputation = {
        "GKP": {"google": "4.6", "ora": "65"},
        "Viera": {"google": "4.7", "ora": "76"},
    }

    output = {}
    for name in ["GKP", "Viera"]:
        units = avail_2026[name]["units"]
        gpr_budget_qtr = income_2026[name]["Gross Potential Rent (GPR)"]["budget"]
        budget_monthly_rent = gpr_budget_qtr / 3 / units if units else 0
        output[name] = {
            "units": units,
            "physical_occ": avail_2026[name]["physical_occ"],
            "leased": avail_2026[name]["leased"],
            "avg_net_eff": avail_2026[name]["avg_net_eff"],
            "avg_sched": avail_2026[name]["avg_sched"],
            "budget_avg_rent": budget_monthly_rent,
            "submarket_avg_rent": submarket_avg_rent[name],
            "available_units": avail_2026[name]["available"],
            "vacant_rented": avail_2026[name]["vacant_rented"],
            "notice_rented": avail_2026[name]["notice_rented"],
            "leads": leads_2026[name]["leads"],
            "tours": leads_2026[name]["tours"],
            "leases": conv_2026[name]["leases_completed"],
            "move_ins": pulse_2026[name]["move_ins"],
            "notices": pulse_2026[name]["notices"],
            "renewals": pulse_2026[name]["renewal_approvals"],
            "income_actual": income_2026[name]["Total Income"]["actual"],
            "income_budget": income_2026[name]["Total Income"]["budget"],
            "income_variance_pct": income_2026[name]["Total Income"]["variance_pct"],
            "nri_actual": income_2026[name]["Net Rental Income"]["actual"],
            "nri_budget": income_2026[name]["Net Rental Income"]["budget"],
            "nri_variance_pct": income_2026[name]["Net Rental Income"]["variance_pct"],
            "prior_physical_occ": avail_2025[name]["physical_occ"],
            "prior_leased": avail_2025[name]["leased"],
            "prior_leads": leads_2025[name]["leads"],
            "prior_tours": leads_2025[name]["tours"],
            "prior_leases": conv_2025[name]["leases_completed"],
            "prior_move_ins": pulse_2025[name]["move_ins"],
            "prior_nri_actual": income_2025[name]["Net Rental Income"]["actual"],
            "prior_income_actual": income_2025[name]["Total Income"]["actual"],
            "google_rating": reputation[name]["google"],
            "ora_score": reputation[name]["ora"],
        }
    return output


def add_rect(slide, left, top, width, height, fill_color, line_color=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=18, color=INK, bold=False,
                font_name="Helvetica", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, items, font_size=16, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.space_after = Pt(4)
        p.font.name = "Helvetica"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
    return box


def add_header(slide, title, subtitle):
    add_rect(slide, 0, 0, Inches(13.333), Inches(0.92), NAVY)
    add_rect(slide, 0, Inches(0.92), Inches(13.333), Inches(0.08), BLUE)
    add_textbox(slide, Inches(0.45), Inches(0.18), Inches(8.2), Inches(0.34), title,
                font_size=26, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.45), Inches(0.52), Inches(10.5), Inches(0.2), subtitle,
                font_size=11, color=RGBColor(218, 234, 240))
    logo_path = ensure_logo_png()
    if logo_path and logo_path.exists():
        add_rect(slide, Inches(11.55), Inches(0.16), Inches(1.28), Inches(0.5), WHITE, line_color=WHITE, radius=True)
        slide.shapes.add_picture(str(logo_path), Inches(11.78), Inches(0.24), width=Inches(0.82))


def add_metric_chip(slide, left, top, width, label, value, fill, accent):
    add_rect(slide, left, top, width, Inches(0.84), fill)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.12), width - Inches(0.36), Inches(0.18),
                label.upper(), font_size=9, color=accent, bold=True)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.32), width - Inches(0.36), Inches(0.3),
                value, font_size=20, color=WHITE, bold=True)


def add_info_card(slide, left, top, width, height, label, value, note="", value_color=INK):
    add_rect(slide, left, top, width, height, CARD)
    add_textbox(
        slide,
        left + Inches(0.18),
        top + Inches(0.12),
        width - Inches(0.36),
        Inches(0.18),
        label,
        font_size=10,
        color=MUTED,
        bold=True,
    )
    add_textbox(
        slide,
        left + Inches(0.18),
        top + Inches(0.34),
        width - Inches(0.36),
        Inches(0.22),
        value,
        font_size=18,
        color=value_color,
        bold=True,
    )
    if note:
        add_textbox(
            slide,
            left + Inches(0.18),
            top + Inches(0.6),
            width - Inches(0.36),
            height - Inches(0.68),
            note,
            font_size=9.5,
            color=INK,
        )


def add_property_panel(slide, left, top, width, title, subtitle, accent, data, yoy_note):
    add_rect(slide, left, top, width, Inches(5.55), WHITE, line_color=BORDER)
    add_rect(slide, left, top, width, Inches(0.86), accent)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.16), width - Inches(0.5), Inches(0.28),
                title, font_size=21, color=WHITE, bold=True)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.49), width - Inches(0.5), Inches(0.2),
                subtitle, font_size=11, color=RGBColor(219, 239, 247),
                bold=True)

    chip_top = top + Inches(0.96)
    chip_gap = Inches(0.1)
    chip_width = (width - Inches(0.4) - 2 * chip_gap) / 3
    add_metric_chip(slide, left + Inches(0.2), chip_top, chip_width, "Physical", pct(data["physical_occ"]), accent, WHITE)
    add_metric_chip(slide, left + Inches(0.2) + chip_width + chip_gap, chip_top, chip_width, "Leased", pct(data["leased"]), accent, WHITE)
    add_metric_chip(slide, left + Inches(0.2) + 2 * (chip_width + chip_gap), chip_top, chip_width, "Q1 Leases", f'{data["leases"]}', accent, WHITE)

    card_top = chip_top + Inches(0.98)
    card_height = Inches(0.7)
    for idx, (label, value) in enumerate([
        ("Leads / Tours / Move-ins", f'{data["leads"]} / {data["tours"]} / {data["move_ins"]}'),
        ("Avg Net Eff. Rent", f'{money(data["avg_net_eff"])} | Budget {money(data["budget_avg_rent"])}'),
        ("Submarket Avg Rent", f'{money(data["submarket_avg_rent"])} | Premium {delta_pct(data["avg_net_eff"], data["submarket_avg_rent"]) or "n/a"}'),
        ("Q1 Total Income", f'{money(data["income_actual"])} vs {money(data["income_budget"])} budget'),
        ("Q1 Net Rental Income", f'{money(data["nri_actual"])} vs {money(data["nri_budget"])} budget'),
    ]):
        y = card_top + Inches(idx * 0.74)
        add_rect(slide, left + Inches(0.2), y, width - Inches(0.4), card_height, CARD)
        add_textbox(slide, left + Inches(0.38), y + Inches(0.12), Inches(2.4), Inches(0.18), label,
                    font_size=11, color=MUTED, bold=True)
        add_textbox(slide, left + Inches(2.95), y + Inches(0.11), width - Inches(3.3), Inches(0.2), value,
                    font_size=14, color=INK, bold=True, align=PP_ALIGN.RIGHT)

    notes = [
        yoy_note,
        f'Available units: {data["available_units"]} | Vacant rented: {data["vacant_rented"]} | Notices rented: {data["notice_rented"]}',
        f'Income variance to budget: {data["income_variance_pct"] * 100:+.1f}% | NRI variance: {data["nri_variance_pct"] * 100:+.1f}%',
    ]
    add_bullets(slide, left + Inches(0.28), top + Inches(4.78), width - Inches(0.56), Inches(0.7), notes, font_size=11, color=INK)


def build_slide_one(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "RISE 55+ Q1 2026 Performance Review",
        "Reporting window: January 1-March 31, 2026. Prior-year operational box score provided only as a January 2025 snapshot.",
    )

    gkp_yoy = "GKP prior-year comparison is not meaningful: Q1 2025 income and occupancy were effectively zero before lease-up activity."
    viera_yoy = (
        f'Viera physical occupancy improved from {pct(data["Viera"]["prior_physical_occ"])} (Jan 31, 2025) '
        f'to {pct(data["Viera"]["physical_occ"])}; Q1 NRI improved from {money(data["Viera"]["prior_nri_actual"])} '
        f'to {money(data["Viera"]["nri_actual"])}.'
    )

    add_property_panel(
        slide,
        Inches(0.35),
        Inches(1.24),
        Inches(6.15),
        "RISE Glen Kernan Park",
        "Jacksonville, FL  |  Early lease-up",
        NAVY,
        data["GKP"],
        gkp_yoy,
    )
    add_property_panel(
        slide,
        Inches(6.83),
        Inches(1.24),
        Inches(6.15),
        "RISE Viera",
        "Viera, FL  |  Lease-up",
        BLUE,
        data["Viera"],
        viera_yoy,
    )

    add_textbox(
        slide,
        Inches(0.45),
        Inches(6.88),
        Inches(12.4),
        Inches(0.22),
        "Source set: Q1 2026 box score exports, Q1 2026 income statement budget-vs-actual workbook, and January 2025 operational snapshot / Q1 2025 financial workbook.",
        font_size=9,
        color=MUTED,
    )


def build_slide_absorption(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Filtered 55+ / Active-Adult Absorption",
        "Absorption proxy = leases in last 30 days / total units from the March 24, 2026 market reports. Quarterly and annual figures are run-rate equivalents.",
    )

    for idx, key in enumerate(["GKP", "Viera"]):
        info = FILTERED_ABSORPTION[key]
        accent = NAVY if key == "GKP" else BLUE
        left = Inches(0.35) if idx == 0 else Inches(6.83)

        add_rect(slide, left, Inches(1.22), Inches(6.15), Inches(5.4), WHITE, line_color=BORDER)
        add_rect(slide, left, Inches(1.22), Inches(6.15), Inches(0.62), accent)
        add_textbox(
            slide,
            left + Inches(0.22),
            Inches(1.4),
            Inches(2.9),
            Inches(0.2),
            info["subject_name"],
            font_size=18,
            color=WHITE,
            bold=True,
        )
        add_rect(
            slide,
            left + Inches(4.3),
            Inches(1.34),
            Inches(1.55),
            Inches(0.28),
            info["status_color"],
            line_color=info["status_color"],
            radius=True,
        )
        add_textbox(
            slide,
            left + Inches(4.34),
            Inches(1.38),
            Inches(1.47),
            Inches(0.12),
            info["status"].upper(),
            font_size=8.5,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

        card_top = Inches(1.98)
        gap = Inches(0.12)
        card_width = Inches(1.8)
        add_info_card(
            slide,
            left + Inches(0.18),
            card_top,
            card_width,
            Inches(1.05),
            "Subject monthly pace",
            pct(info["subject_abs"]),
            f'{info["subject_leases30"]} leases / {info["subject_units"]} units',
            value_color=accent,
        )
        add_info_card(
            slide,
            left + Inches(0.18) + card_width + gap,
            card_top,
            card_width,
            Inches(1.05),
            "Filtered submarket",
            pct(info["submarket_abs"]),
            f'Quarterly {pct(info["submarket_abs"] * 3)} | Annualized {pct(info["submarket_abs"] * 12)}',
        )
        add_info_card(
            slide,
            left + Inches(0.18) + 2 * (card_width + gap),
            card_top,
            card_width,
            Inches(1.05),
            "Leased position",
            f'{pct(info["subject_leased"])} vs {pct(info["submarket_leased"])}',
            f'Gap to filtered peer set: {info["delta_leased"] * 100:+.1f} pts',
        )

        add_rect(slide, left + Inches(0.18), Inches(3.18), Inches(2.6), Inches(1.18), CARD)
        add_textbox(
            slide,
            left + Inches(0.38),
            Inches(3.32),
            Inches(2.2),
            Inches(0.18),
            "Filtered comp set",
            font_size=11,
            color=MUTED,
            bold=True,
        )
        add_bullets(
            slide,
            left + Inches(0.33),
            Inches(3.56),
            Inches(2.22),
            Inches(0.66),
            info["comps"],
            font_size=10.5,
            color=INK,
        )

        add_rect(slide, left + Inches(2.92), Inches(3.18), Inches(3.05), Inches(1.18), CARD)
        add_textbox(
            slide,
            left + Inches(3.12),
            Inches(3.32),
            Inches(2.65),
            Inches(0.18),
            "Pace read",
            font_size=11,
            color=MUTED,
            bold=True,
        )
        add_textbox(
            slide,
            left + Inches(3.12),
            Inches(3.55),
            Inches(2.58),
            Inches(0.28),
            f'Pace variance: {info["delta_abs"] * 100:+.1f} pts',
            font_size=15,
            color=info["status_color"] if info["delta_abs"] >= 0 else AMBER,
            bold=True,
        )
        add_textbox(
            slide,
            left + Inches(3.12),
            Inches(3.84),
            Inches(2.58),
            Inches(0.36),
            info["growth_note"],
            font_size=9.5,
            color=INK,
        )

        add_rect(slide, left + Inches(0.18), Inches(4.55), Inches(5.79), Inches(1.72), NAVY)
        add_textbox(
            slide,
            left + Inches(0.38),
            Inches(4.72),
            Inches(5.4),
            Inches(0.18),
            info["market"],
            font_size=12,
            color=WHITE,
            bold=True,
        )
        summary = [
            f'{info["subject_name"]} is {info["status"].lower()} versus the filtered submarket on current leasing velocity.',
            f'Current leased percentage remains below the filtered benchmark by {abs(info["delta_leased"]) * 100:.1f} points.',
        ]
        add_bullets(
            slide,
            left + Inches(0.34),
            Inches(5.0),
            Inches(5.2),
            Inches(0.78),
            summary,
            font_size=11,
            color=WHITE,
        )

    add_textbox(
        slide,
        Inches(0.45),
        Inches(6.9),
        Inches(12.35),
        Inches(0.16),
        "Filtered benchmark excludes conventional multifamily and non-direct senior living where the shared report suggested a weaker apples-to-apples comparison. Viera remains directional because the filtered set includes one clean active-adult comp.",
        font_size=8.8,
        color=MUTED,
    )


def build_slide_two(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "30 / 60 / 90 Day Plan",
        "30-day actions now include the confirmed in-person training schedule and the follow-up roadmap for Morgan Matlock support.",
    )

    plans = {
        "GKP": {
            "title": "GKP",
            "accent": NAVY,
            "current": [
                "Market analysis and marketing plan adjustments in motion: SEO, ads, mailers",
                "Pricing reset completed",
                "Leasing Manager hired",
                "Fast-fill bonus structure launched to drive velocity",
                "All-in pricing rolled out",
                f'{GKP_TRAINING_DATE} in-person training: lead intake, first response, discovery, follow-up, and closing execution',
            ],
            "next": [
                "Hire Leasing Professional",
                "Hire Lifestyle Coordinator",
                "Update discovery process and add testimonials into the tour path",
                "Continue submarket checks to balance velocity with asset value protection",
                "Prepare first renewal offers and referral-program launch once occupancy exceeds 25%",
            ],
            "later": [
                "Monitor vacant utilities and maintenance costs closely while occupancy builds",
                "Evaluate guest suite / model use once occupancy reaches 40%-50%",
            ],
        },
        "Viera": {
            "title": "Viera",
            "accent": BLUE,
            "current": [
                "General Manager hired",
                "Leasing Manager hired",
                "All-in pricing rolled out",
                f'{VIERA_TRAINING_DATE} in-person training: lead intake, first response, discovery, follow-up, and closing execution',
            ],
            "next": [
                "Refresh market analysis and marketing plan: SEO, ads, mailers",
                "Deploy aggressive fast-fill bonus structure",
                "Tighten discovery process and add testimonials",
                "Push renewal conversations and resident referrals as occupancy grows",
            ],
            "later": [
                "Protect rate while improving qualified-tour conversion",
                "Expand review mix beyond Google to improve ORA visibility",
            ],
        },
    }

    for idx, key in enumerate(["GKP", "Viera"]):
        left = Inches(0.35) if idx == 0 else Inches(6.83)
        accent = plans[key]["accent"]
        add_rect(slide, left, Inches(1.2), Inches(6.15), Inches(5.0), WHITE, line_color=BORDER)
        add_rect(slide, left, Inches(1.2), Inches(6.15), Inches(0.08), accent)
        add_rect(slide, left + Inches(0.18), Inches(1.33), Inches(4.1), Inches(0.34), accent)
        add_textbox(slide, left + Inches(0.18), Inches(1.35), Inches(4.1), Inches(0.2), plans[key]["title"],
                    font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        sections = [
            ("30 Days", plans[key]["current"], Inches(1.78), Inches(1.12)),
            ("31-60 Days", plans[key]["next"], Inches(2.98), Inches(1.75)),
            ("61-90 Days", plans[key]["later"], Inches(4.83), Inches(0.92)),
        ]
        for section_title, items, top, height in sections:
            add_rect(slide, left + Inches(0.18), top, Inches(5.79), height, CARD)
            add_textbox(slide, left + Inches(0.35), top + Inches(0.1), Inches(2.2), Inches(0.18), section_title,
                        font_size=11, color=MUTED, bold=True)
            add_bullets(slide, left + Inches(0.28), top + Inches(0.28), Inches(5.48), height - Inches(0.24), items,
                        font_size=10 if key == "GKP" else 11, color=INK)

    add_rect(slide, Inches(0.35), Inches(6.34), Inches(12.63), Inches(0.68), NAVY)
    training = [
        f'Viera in-person {VIERA_TRAINING_DATE} | GKP in-person {GKP_TRAINING_DATE}',
        f'May follow-up with Morgan Matlock: Viera {VIERA_FOLLOWUP_DATE} | GKP {GKP_FOLLOWUP_DATE}',
        "Ongoing virtual follow-up calls, corporate support, sentiment review, and weekly coaching when benchmarks are missed",
    ]
    add_textbox(slide, Inches(0.58), Inches(6.47), Inches(2.2), Inches(0.2), "Training Roadmap",
                font_size=12, color=WHITE, bold=True)
    add_textbox(slide, Inches(2.9), Inches(6.45), Inches(9.8), Inches(0.22), "  |  ".join(training),
                font_size=10.5, color=WHITE)


def build_slide_three(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Marketing Strategy + Next 90 Days",
        "Focus: keep Viera at or above budget while accelerating Glen Kernan Park lease-up without giving away the rate premium to submarket.",
    )

    add_rect(slide, Inches(0.35), Inches(1.25), Inches(4.05), Inches(5.45), WHITE, line_color=BORDER)
    add_rect(slide, Inches(0.35), Inches(1.25), Inches(4.05), Inches(0.54), NAVY)
    add_textbox(slide, Inches(0.58), Inches(1.4), Inches(3.3), Inches(0.2), "Current Marketing", font_size=18, color=WHITE, bold=True)
    current_items = [
        "SEO / paid search, digital ads, and direct mail in market",
        "Pricing reset and all-in pricing messaging already rolled out",
        "Fast-fill bonus plan to improve lease velocity on exposed units",
        "Discovery-process updates and testimonial integration underway",
        "Weekly submarket review to stay competitive without eroding value",
        "Review / reputation push: GKP Google 4.6 | ORA 65; Viera Google 4.7 | ORA 76",
    ]
    add_bullets(slide, Inches(0.56), Inches(1.96), Inches(3.62), Inches(3.45), current_items, font_size=12, color=INK)

    add_rect(slide, Inches(0.55), Inches(5.53), Inches(3.64), Inches(0.95), CARD)
    add_textbox(slide, Inches(0.75), Inches(5.68), Inches(3.2), Inches(0.18), "Marketing takeaway", font_size=11, color=MUTED, bold=True)
    add_textbox(
        slide,
        Inches(0.75),
        Inches(5.92),
        Inches(3.05),
        Inches(0.36),
        "GKP still has a strong rate premium to submarket; Viera is near budget on revenue, so the next lift should come from conversion and retention before broader price pressure.",
        font_size=10.5,
        color=INK,
    )

    for idx, key in enumerate(["GKP", "Viera"]):
        left = Inches(4.68) if idx == 0 else Inches(8.82)
        accent = NAVY if key == "GKP" else BLUE
        add_rect(slide, left, Inches(1.25), Inches(4.1), Inches(5.45), WHITE, line_color=BORDER)
        add_rect(slide, left, Inches(1.25), Inches(4.1), Inches(0.54), accent)
        add_textbox(slide, left + Inches(0.22), Inches(1.4), Inches(2.5), Inches(0.2), key, font_size=18, color=WHITE, bold=True)

        if key == "GKP":
            next_items = [
                "Concentrate ad spend and pricing action on hardest-to-lease plans and highest-exposure homes",
                "Use local 55+ partnerships, broker outreach, and events to add qualified tours",
                "Close the gap between 14.3% leased and budget with weekly follow-up accountability",
                "Maintain rate discipline: current effective rent is materially above both budget and submarket",
                "Cross the occupancy threshold needed to activate a resident-referral program",
            ]
            goal_line = "90-day focus metrics: increase tour volume, close more than 7 leases per month, and narrow the budget revenue gap."
        else:
            next_items = [
                "Protect NRI that is already slightly above budget while improving occupancy",
                "Tighten lead-to-tour and tour-to-lease conversion through scripting and follow-up cadence",
                "Lean into referrals, renewals, and resident reviews to lower acquisition cost",
                "Diversify review sources beyond Google so ORA improves with volume",
                "Keep rate above submarket while avoiding unnecessary concession drift",
            ]
            goal_line = "90-day focus metrics: push physical occupancy into the upper-50s / low-60s while holding NRI at or above budget."

        add_bullets(slide, left + Inches(0.22), Inches(1.94), Inches(3.66), Inches(3.2), next_items, font_size=11, color=INK)
        add_rect(slide, left + Inches(0.18), Inches(5.53), Inches(3.74), Inches(0.95), CARD)
        add_textbox(slide, left + Inches(0.35), Inches(5.68), Inches(3.38), Inches(0.4), goal_line, font_size=10.5, color=INK, bold=True)

    add_textbox(
        slide,
        Inches(0.45),
        Inches(6.86),
        Inches(12.4),
        Inches(0.18),
        f'Data note: Viera prior-year operating comparison uses the January 31, 2025 snapshot that was provided; GKP prior-year operations were pre-opening. Deck generated {datetime.now().strftime("%B %d, %Y")}.',
        font_size=8.8,
        color=MUTED,
    )


def build_deck():
    data = assemble_data()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_one(prs, data)
    build_slide_absorption(prs)
    build_slide_two(prs)
    build_slide_three(prs, data)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build_deck()
    print(path)
